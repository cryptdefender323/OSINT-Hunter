#!/usr/bin/env python3

import os
import json
import hashlib
import re
import sys
from datetime import datetime
from rich.console import Console
from rich.table import Table
from rich.panel import Panel

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from config import Config

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None
try:
    import docx
except ImportError:
    docx = None
try:
    import exifread
except ImportError:
    exifread = None
try:
    from PIL import Image
except ImportError:
    Image = None
try:
    from mutagen import File as MutagenFile
except ImportError:
    MutagenFile = None
try:
    import openpyxl
except ImportError:
    openpyxl = None
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

console = Console()

SENSITIVE_PATTERNS = {
    "email": r"[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}",
    "phone_id": r"(?:\+62|62|0)(?:\d{9,12})",
    "phone_intl": r"\+?\d{1,4}[\s-]?\(?\d{1,4}\)?[\s-]?\d{3,4}[\s-]?\d{3,4}",
    "credit_card": r"\b(?:\d{4}[\s-]?){3}\d{4}\b",
    "ipv4": r"\b(?:\d{1,3}\.){3}\d{1,3}\b",
    "api_key": r"(?:api[_-]?key|apikey|access[_-]?token|secret[_-]?key)\s*[:=]\s*['\"]?([a-zA-Z0-9_\-]{20,})['\"]?",
    "jwt": r"eyJ[a-zA-Z0-9_-]+\.eyJ[a-zA-Z0-9_-]+\.[a-zA-Z0-9_-]+",
    "password": r"(?:password|passwd|pwd)\s*[:=]\s*['\"]?([^\s'\"]+)",
    "nik_ktp": r"\b\d{16}\b",
    "url": r"https?://[^\s<>\"']+",
    "aws_key": r"AKIA[0-9A-Z]{16}",
    "private_key": r"-----BEGIN\s+(?:RSA|EC|DSA|OPENSSH)?\s*PRIVATE\s+KEY-----",
}


class MetadataExtractor:
    def __init__(self, filepath=""):
        self.filepath = filepath

    def compute_hashes(self):
        hashes = {}
        try:
            with open(self.filepath, "rb") as f:
                content = f.read()
            hashes["md5"] = hashlib.md5(content).hexdigest()
            hashes["sha1"] = hashlib.sha1(content).hexdigest()
            hashes["sha256"] = hashlib.sha256(content).hexdigest()
        except Exception as e:
            hashes["error"] = str(e)
        return hashes

    def get_file_info(self):
        try:
            stat = os.stat(self.filepath)
            size = stat.st_size
            if size < 1024:
                size_str = f"{size} B"
            elif size < 1048576:
                size_str = f"{size / 1024:.1f} KB"
            else:
                size_str = f"{size / 1048576:.1f} MB"

            return {
                "filename": os.path.basename(self.filepath),
                "full_path": os.path.abspath(self.filepath),
                "size_bytes": size,
                "size_human": size_str,
                "created": datetime.fromtimestamp(stat.st_ctime).isoformat(),
                "modified": datetime.fromtimestamp(stat.st_mtime).isoformat(),
                "extension": os.path.splitext(self.filepath)[-1].lower(),
                "hashes": self.compute_hashes(),
            }
        except Exception as e:
            return {"error": str(e)}

    def extract_pdf(self):
        if not PyPDF2:
            return {"error": "PyPDF2 not installed"}
        data = {}
        try:
            with open(self.filepath, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                if reader.metadata:
                    data["metadata"] = {k: str(v) for k, v in reader.metadata.items() if v}
                data["page_count"] = len(reader.pages)
                data["is_encrypted"] = reader.is_encrypted
                text_pages = []
                for i, page in enumerate(reader.pages):
                    text = page.extract_text()
                    if text and text.strip():
                        text_pages.append({"page": i + 1, "text": text.strip()[:500]})
                data["text_preview"] = text_pages[:5]
        except Exception as e:
            data["error"] = str(e)
        return data

    def extract_docx(self):
        if not docx:
            return {"error": "python-docx not installed"}
        data = {}
        try:
            doc = docx.Document(self.filepath)
            props = doc.core_properties
            data["properties"] = {
                "author": props.author,
                "title": props.title,
                "subject": props.subject,
                "created": str(props.created),
                "modified": str(props.modified),
                "last_modified_by": props.last_modified_by,
                "revision": props.revision,
                "category": props.category,
                "comments": props.comments,
                "keywords": props.keywords,
            }
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            data["paragraph_count"] = len(paragraphs)
            data["text_preview"] = paragraphs[:10]
            data["table_count"] = len(doc.tables)
            data["image_count"] = len(doc.inline_shapes)
        except Exception as e:
            data["error"] = str(e)
        return data

    def extract_xlsx(self):
        if not openpyxl:
            return {"error": "openpyxl not installed"}
        data = {}
        try:
            wb = openpyxl.load_workbook(self.filepath, read_only=True, data_only=True)
            data["sheet_names"] = wb.sheetnames
            data["sheet_count"] = len(wb.sheetnames)
            data["properties"] = {
                "creator": wb.properties.creator,
                "title": wb.properties.title,
                "created": str(wb.properties.created),
                "modified": str(wb.properties.modified),
                "last_modified_by": wb.properties.lastModifiedBy,
            }
            sheets_info = []
            for name in wb.sheetnames[:3]:
                ws = wb[name]
                sheets_info.append({
                    "name": name,
                    "dimensions": ws.dimensions,
                    "max_row": ws.max_row,
                    "max_column": ws.max_column,
                })
            data["sheets"] = sheets_info
            wb.close()
        except Exception as e:
            data["error"] = str(e)
        return data

    def extract_pptx(self):
        if not Presentation:
            return {"error": "python-pptx not installed"}
        data = {}
        try:
            prs = Presentation(self.filepath)
            data["slide_count"] = len(prs.slides)
            data["slide_width"] = str(prs.slide_width)
            data["slide_height"] = str(prs.slide_height)
            slides_info = []
            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text[:100])
                slides_info.append({"slide": i + 1, "texts": texts[:5]})
            data["slides_preview"] = slides_info[:5]
        except Exception as e:
            data["error"] = str(e)
        return data

    def extract_image_exif(self):
        if not exifread:
            return {"error": "exifread not installed"}
        exif_data = {}
        try:
            with open(self.filepath, 'rb') as f:
                tags = exifread.process_file(f)
            for tag in tags:
                exif_data[str(tag)] = str(tags[tag])

            if 'GPS GPSLatitude' in tags and 'GPS GPSLongitude' in tags:
                gps = self._parse_gps(tags)
                if gps:
                    exif_data['gps_coordinates'] = gps
        except Exception as e:
            exif_data["error"] = str(e)
        return exif_data

    def _parse_gps(self, tags):
        try:
            def to_degrees(val):
                d, m, s = [float(x.num) / float(x.den) for x in val.values]
                return d + (m / 60.0) + (s / 3600.0)

            lat = to_degrees(tags['GPS GPSLatitude'])
            lon = to_degrees(tags['GPS GPSLongitude'])
            if str(tags.get('GPS GPSLatitudeRef')) == 'S':
                lat = -lat
            if str(tags.get('GPS GPSLongitudeRef')) == 'W':
                lon = -lon
            return {
                "latitude": lat,
                "longitude": lon,
                "google_maps": f"https://maps.google.com/?q={lat},{lon}"
            }
        except Exception:
            return None

    def extract_audio_video(self):
        if not MutagenFile:
            return {"error": "mutagen not installed"}
        data = {}
        try:
            audio = MutagenFile(self.filepath)
            if audio:
                data["format"] = type(audio).__name__
                data["length_seconds"] = round(audio.info.length, 2) if hasattr(audio.info, 'length') else None
                data["bitrate"] = audio.info.bitrate if hasattr(audio.info, 'bitrate') else None
                data["sample_rate"] = audio.info.sample_rate if hasattr(audio.info, 'sample_rate') else None
                data["channels"] = audio.info.channels if hasattr(audio.info, 'channels') else None
                if hasattr(audio, 'tags') and audio.tags:
                    tags = {}
                    for key in audio.tags:
                        try:
                            tags[str(key)] = str(audio.tags[key])[:100]
                        except Exception:
                            pass
                    data["tags"] = tags
        except Exception as e:
            data["error"] = str(e)
        return data

    def scan_sensitive_data(self, text_content):
        findings = []
        for category, pattern in SENSITIVE_PATTERNS.items():
            matches = re.findall(pattern, text_content, re.IGNORECASE)
            if matches:
                unique_matches = list(set(matches))[:10]
                findings.append({
                    "category": category,
                    "count": len(matches),
                    "samples": unique_matches[:5],
                })
        return findings

    def extract_all_text(self):
        ext = os.path.splitext(self.filepath)[-1].lower()
        text = ""
        if ext == ".pdf" and PyPDF2:
            try:
                with open(self.filepath, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        t = page.extract_text()
                        if t:
                            text += t + "\n"
            except Exception:
                pass
        elif ext == ".docx" and docx:
            try:
                doc = docx.Document(self.filepath)
                text = "\n".join(p.text for p in doc.paragraphs)
            except Exception:
                pass
        elif ext == ".txt":
            try:
                with open(self.filepath, "r", errors="ignore") as f:
                    text = f.read()
            except Exception:
                pass
        return text

    def extract_metadata(self):
        ext = os.path.splitext(self.filepath)[-1].lower()
        file_info = self.get_file_info()
        if "error" in file_info:
            return file_info

        result = {"file_info": file_info, "type_data": {}, "sensitive_findings": []}

        if ext == ".pdf":
            result["type_data"] = self.extract_pdf()
        elif ext == ".docx":
            result["type_data"] = self.extract_docx()
        elif ext == ".xlsx":
            result["type_data"] = self.extract_xlsx()
        elif ext == ".pptx":
            result["type_data"] = self.extract_pptx()
        elif ext in [".jpg", ".jpeg", ".png", ".tiff", ".bmp"]:
            result["type_data"] = self.extract_image_exif()
        elif ext in [".mp3", ".mp4", ".m4a", ".flac", ".ogg", ".wav", ".avi", ".mkv"]:
            result["type_data"] = self.extract_audio_video()
        elif ext == ".txt":
            result["type_data"] = {"content_preview": open(self.filepath, "r", errors="ignore").read()[:500]}
        else:
            result["type_data"] = {"note": f"Basic file info extracted. Extension '{ext}' has no specialized parser."}

        all_text = self.extract_all_text()
        if all_text:
            result["sensitive_findings"] = self.scan_sensitive_data(all_text)

        return result

    def run(self):
        console.print(Panel(
            "[bold cyan]METADATA EXTRACTOR — OSINT-Hunter V3[/bold cyan]\n"
            "[dim]PDF • DOCX • XLSX • PPTX • Images • Audio/Video • Sensitive Data Scanner[/dim]",
            border_style="cyan"
        ))

        filepath = input("\n  Enter file path: ").strip().strip("'\"")
        self.filepath = filepath

        if not os.path.exists(self.filepath):
            console.print("[red]❌ File not found![/red]")
            return

        console.print(f"[blue]→ Analyzing: {self.filepath}[/blue]")
        result = self.extract_metadata()

        fi = result.get("file_info", {})
        info_table = Table(title="File Information", show_header=False)
        info_table.add_column("Field", style="cyan", width=15)
        info_table.add_column("Value", style="white")
        for k, v in fi.items():
            if k != "hashes":
                info_table.add_row(k.replace("_", " ").title(), str(v))
        if "hashes" in fi:
            for algo, h in fi["hashes"].items():
                info_table.add_row(f"Hash ({algo.upper()})", str(h))
        console.print(info_table)

        td = result.get("type_data", {})
        if td and "error" not in td:
            type_table = Table(title="Extracted Metadata", show_header=False)
            type_table.add_column("Field", style="cyan", width=20)
            type_table.add_column("Value", style="yellow")
            for k, v in td.items():
                if not isinstance(v, (list, dict)):
                    type_table.add_row(k.replace("_", " ").title(), str(v)[:100])
                elif isinstance(v, dict):
                    for sk, sv in v.items():
                        type_table.add_row(f"  {sk}", str(sv)[:100])
            console.print(type_table)

        findings = result.get("sensitive_findings", [])
        if findings:
            console.print(Panel("[bold red]⚠ SENSITIVE DATA DETECTED[/bold red]", border_style="red"))
            sens_table = Table(title="Sensitive Data Findings", show_header=True, header_style="bold red")
            sens_table.add_column("Category", style="cyan")
            sens_table.add_column("Count", justify="right", style="magenta")
            sens_table.add_column("Samples", style="yellow")
            for f in findings:
                sens_table.add_row(f["category"].upper(), str(f["count"]), ", ".join(str(s) for s in f["samples"][:3]))
            console.print(sens_table)

        Config.ensure_dirs()
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        outname = f"logs/metadata_{os.path.basename(self.filepath)}_{timestamp}.json"
        with open(outname, "w") as f:
            json.dump(result, f, indent=2, ensure_ascii=False, default=str)
        console.print(f"\n[green]✔ Full results saved to {outname}[/green]")


def main():
    extractor = MetadataExtractor()
    extractor.run()


if __name__ == "__main__":
    main()
